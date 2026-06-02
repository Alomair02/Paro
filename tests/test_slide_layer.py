import tempfile
import unittest
from pathlib import Path
from zipfile import ZipFile

from lxml import etree

from builders.common import REFERENCE
from builders.deck_builder import build_deck
from builders.shape_emitters import (
    SlideState,
    emit_autoshape,
    emit_image,
    emit_line,
    emit_placeholder_text,
    emit_slide_background,
    emit_table,
    emit_text_box,
)
from builders.slide_builder import SlideBuilder
from builders.zip_assembler import ZIPAssembler
from core.content_type_reg import ContentTypeRegistry
from core.relationship_reg import RelationshipRegistry
from core.xml_builder import CANONICAL_XML_DECLARATION, NSMAP, qn
from tests.pptx_test_utils import CT_NS, REL_NS, build_core_parts, parse_xml


PNG_BYTES = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
    b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00"
    b"\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00"
    b"\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)


def rels_for(registry: RelationshipRegistry, part_path: str):
    return parse_xml(registry.render(part_path)).findall("rel:Relationship", REL_NS)


def c_nv_pr_ids(root: etree._Element) -> list[str]:
    return [
        c_nv_pr.get("id")
        for c_nv_pr in root.findall(".//p:cNvPr", NSMAP)
        if c_nv_pr.get("id")
    ]


def z_order_names(root: etree._Element) -> list[str]:
    sp_tree = root.find("p:cSld/p:spTree", NSMAP)
    names = []
    for child in list(sp_tree)[2:]:
        c_nv_pr = child.find("p:nvSpPr/p:cNvPr", NSMAP)
        if c_nv_pr is None:
            c_nv_pr = child.find("p:nvPicPr/p:cNvPr", NSMAP)
        names.append(c_nv_pr.get("name"))
    return names


class ShapeEmitterTests(unittest.TestCase):
    def test_placeholder_without_geometry_omits_xfrm(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_placeholder_text(
            {
                "idx": 0,
                "placeholder_type": "ctrTitle",
                "name": "Title Placeholder",
                "text": "Inherited Title",
            },
            state,
        )

        self.assertEqual(shape.tag, qn("p", "sp"))
        self.assertEqual(shape.find("p:nvSpPr/p:cNvPr", NSMAP).get("id"), "2")
        ph = shape.find("p:nvSpPr/p:nvPr/p:ph", NSMAP)
        self.assertEqual(ph.get("idx"), "0")
        self.assertEqual(ph.get("type"), "ctrTitle")
        self.assertIsNone(shape.find("p:spPr/a:xfrm", NSMAP))
        self.assertEqual(shape.find(".//a:t", NSMAP).text, "Inherited Title")

    def test_placeholder_with_explicit_geometry_emits_xfrm_only_then(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_placeholder_text(
            {
                "idx": 1,
                "placeholder_type": "subTitle",
                "x": "1in",
                "y": "0.5in",
                "w": "4in",
                "h": "1in",
                "text": "Subtitle",
            },
            state,
        )

        off = shape.find("p:spPr/a:xfrm/a:off", NSMAP)
        ext = shape.find("p:spPr/a:xfrm/a:ext", NSMAP)
        self.assertEqual(off.get("x"), "914400")
        self.assertEqual(off.get("y"), "457200")
        self.assertEqual(ext.get("cx"), "3657600")
        self.assertEqual(ext.get("cy"), "914400")

    def test_text_box_is_self_contained_and_assigns_id(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_text_box(
            {
                "name": "Body Text Box",
                "x": "1in",
                "y": "2in",
                "w": "3in",
                "h": "1in",
                "text": "Box",
            },
            state,
        )

        self.assertEqual(shape.tag, qn("p", "sp"))
        self.assertEqual(shape.find("p:nvSpPr/p:cNvPr", NSMAP).get("id"), "2")
        self.assertEqual(shape.find("p:nvSpPr/p:cNvSpPr", NSMAP).get("txBox"), "1")
        self.assertIsNotNone(shape.find("p:spPr/a:xfrm", NSMAP))
        self.assertEqual(shape.find("p:spPr/a:prstGeom", NSMAP).get("prst"), "rect")
        self.assertIsNotNone(shape.find("p:spPr/a:noFill", NSMAP))
        self.assertIsNotNone(shape.find("p:spPr/a:ln/a:noFill", NSMAP))

    def test_autoshape_uses_preset_fill_line_and_optional_text(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_autoshape(
            {
                "name": "Accent Rectangle",
                "x": "1in",
                "y": "1in",
                "w": "2in",
                "h": "1in",
                "preset": "roundRect",
                "fill": "accent1",
                "line": False,
                "text": "Accent",
            },
            state,
        )

        self.assertEqual(shape.find("p:spPr/a:prstGeom", NSMAP).get("prst"), "roundRect")
        self.assertEqual(
            shape.find("p:spPr/a:solidFill/a:schemeClr", NSMAP).get("val"),
            "accent1",
        )
        self.assertIsNotNone(shape.find("p:spPr/a:ln/a:noFill", NSMAP))
        self.assertEqual(shape.find(".//a:t", NSMAP).text, "Accent")

    def test_autoshape_roundrect_radius_fill_effects_emit_real_sppr(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_autoshape(
            {
                "name": "Finished Bar",
                "x": "1in",
                "y": "1in",
                "w": "2in",
                "h": "0.3in",
                "preset": "roundRect",
                "radius": "3pt",
                "fill": "accent1",
                "fill_style": {
                    "type": "gradient",
                    "color": "accent1",
                    "startTransforms": {"lumMod": 100000, "lumOff": 18000},
                    "endTransforms": {"lumMod": 82000},
                },
                "line": {"color": "accent1", "width": "0.75pt"},
                "effects": {"shadow": {"blur": "3pt", "dist": "1pt", "alpha": 12000}},
            },
            state,
        )

        gd = shape.find("p:spPr/a:prstGeom/a:avLst/a:gd", NSMAP)
        self.assertIsNotNone(gd)
        self.assertTrue(gd.get("fmla").startswith("val "))
        self.assertIsNotNone(shape.find("p:spPr/a:gradFill", NSMAP))
        self.assertEqual(shape.find("p:spPr/a:ln", NSMAP).get("w"), str(int(REFERENCE["units"]["emu_per"]["pt"] * 0.75)))
        self.assertIsNotNone(shape.find("p:spPr/a:effectLst/a:outerShdw", NSMAP))

    def test_autoshape_fill_none_emits_no_fill(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_autoshape(
            {
                "x": "1in",
                "y": "1in",
                "w": "2in",
                "h": "1in",
                "fill": "none",
                "line": {"color": "accent1", "width": "1pt"},
            },
            state,
        )

        self.assertIsNotNone(shape.find("p:spPr/a:noFill", NSMAP))
        self.assertIsNone(shape.find("p:spPr/a:solidFill", NSMAP))

    def test_text_box_emits_bullets_numbering_underline_and_hyperlink(self):
        state = SlideState("ppt/slides/slide1.xml")
        relationships = RelationshipRegistry()

        shape = emit_text_box(
            {
                "x": "1in",
                "y": "1in",
                "w": "4in",
                "h": "2in",
                "paragraphs": [
                    {"level": 1, "bullet": "bullet", "runs": [{"text": "Bullet"}]},
                    {"bullet": "number", "runs": [{"text": "Number"}]},
                    {
                        "bullet": "none",
                        "runs": [
                            {
                                "text": "Linked",
                                "underline": True,
                                "link": "https://example.com/?a=1&b=2",
                            }
                        ],
                    },
                ],
            },
            state,
            relationships,
        )

        paragraphs = shape.findall("p:txBody/a:p", NSMAP)
        bullet = paragraphs[0].find("a:pPr/a:buChar", NSMAP)
        numbering = paragraphs[1].find("a:pPr/a:buAutoNum", NSMAP)
        no_bullet = paragraphs[2].find("a:pPr/a:buNone", NSMAP)
        self.assertEqual(bullet.get("char"), REFERENCE["text_defaults"]["list_markers"]["bullet"])
        self.assertEqual(numbering.get("type"), REFERENCE["text_defaults"]["numbering_type"])
        self.assertIsNotNone(no_bullet)
        self.assertEqual(
            paragraphs[0].find("a:pPr", NSMAP).get("marL"),
            str(REFERENCE["text_defaults"]["level_indents"][1]["marL"]),
        )
        rpr = paragraphs[2].find("a:r/a:rPr", NSMAP)
        self.assertEqual(rpr.get("u"), "sng")
        self.assertEqual(rpr.find("a:hlinkClick", NSMAP).get(qn("r", "id")), "rId1")

        rel = rels_for(relationships, "ppt/slides/slide1.xml")[0]
        self.assertEqual(rel.get("Type"), RelationshipRegistry.HYPERLINK)
        self.assertEqual(rel.get("Target"), "https://example.com/?a=1&b=2")
        self.assertEqual(rel.get("TargetMode"), "External")

    def test_text_box_emits_paragraph_spacing_from_role_resolution(self):
        state = SlideState("ppt/slides/slide1.xml")

        shape = emit_text_box(
            {
                "x": "1in",
                "y": "1in",
                "w": "4in",
                "h": "1in",
                "paragraphs": [
                    {
                        "runs": [{"text": "Spaced"}],
                        "lineSpacing": 1.2,
                        "spaceBefore": "3pt",
                        "spaceAfter": "6pt",
                    }
                ],
            },
            state,
        )

        p_pr = shape.find("p:txBody/a:p/a:pPr", NSMAP)
        self.assertEqual(p_pr.find("a:lnSpc/a:spcPct", NSMAP).get("val"), "120000")
        self.assertEqual(p_pr.find("a:spcBef/a:spcPts", NSMAP).get("val"), "300")
        self.assertEqual(p_pr.find("a:spcAft/a:spcPts", NSMAP).get("val"), "600")

    def test_line_emits_native_connector_dash_and_cap(self):
        state = SlideState("ppt/slides/slide1.xml")

        line = emit_line(
            {
                "x1": "1in",
                "y1": "1in",
                "x2": "4in",
                "y2": "2in",
                "color": "accent2",
                "width": "2pt",
                "dash": "dash",
                "cap": "round",
            },
            state,
        )

        self.assertEqual(line.tag, qn("p", "cxnSp"))
        self.assertEqual(line.find("p:spPr/a:prstGeom", NSMAP).get("prst"), "line")
        ln = line.find("p:spPr/a:ln", NSMAP)
        self.assertEqual(ln.get("cap"), "rnd")
        self.assertEqual(ln.find("a:solidFill/a:schemeClr", NSMAP).get("val"), "accent2")
        self.assertEqual(ln.find("a:prstDash", NSMAP).get("val"), "dash")

    def test_line_connector_extents_are_never_zero(self):
        state = SlideState("ppt/slides/slide1.xml")

        vertical = emit_line(
            {
                "x1": "1in",
                "y1": "1in",
                "x2": "1in",
                "y2": "3in",
            },
            state,
        )
        horizontal = emit_line(
            {
                "x1": "1in",
                "y1": "1in",
                "x2": "3in",
                "y2": "1in",
            },
            state,
        )

        vertical_ext = vertical.find("p:spPr/a:xfrm/a:ext", NSMAP)
        horizontal_ext = horizontal.find("p:spPr/a:xfrm/a:ext", NSMAP)
        self.assertEqual(vertical_ext.get("cx"), "1")
        self.assertEqual(vertical_ext.get("cy"), str(REFERENCE["units"]["emu_per"]["inch"] * 2))
        self.assertEqual(horizontal_ext.get("cx"), str(REFERENCE["units"]["emu_per"]["inch"] * 2))
        self.assertEqual(horizontal_ext.get("cy"), "1")

    def test_table_emits_native_graphic_frame_borders_and_merges(self):
        state = SlideState("ppt/slides/slide1.xml")

        table = emit_table(
            {
                "x": "1in",
                "y": "1in",
                "w": "4in",
                "h": "1in",
                "columns": [{"w": "2in"}, {"w": "2in"}],
                "header": True,
                "rows": [
                    {
                        "h": "1in",
                        "cells": [
                            {
                                "gridSpan": 2,
                                "rowSpan": 2,
                                "text": "Header",
                                "align": "ctr",
                                "line": {"color": "dk2", "width": "0.75pt"},
                            }
                        ],
                    }
                ],
            },
            state,
        )

        self.assertEqual(table.tag, qn("p", "graphicFrame"))
        self.assertIsNotNone(table.find("a:graphic/a:graphicData/a:tbl", NSMAP))
        self.assertEqual(len(table.findall(".//a:tblGrid/a:gridCol", NSMAP)), 2)
        cell = table.find(".//a:tc", NSMAP)
        self.assertEqual(cell.get("gridSpan"), "2")
        self.assertEqual(cell.get("rowSpan"), "2")
        self.assertEqual(cell.find("a:txBody/a:p/a:r/a:t", NSMAP).text, "Header")
        self.assertIsNotNone(cell.find("a:tcPr/a:solidFill", NSMAP))
        self.assertIsNotNone(cell.find("a:tcPr/a:lnL/a:solidFill", NSMAP))

    def test_image_registers_media_relationship_and_content_type(self):
        content_types = ContentTypeRegistry()
        relationships = RelationshipRegistry()

        with tempfile.TemporaryDirectory() as tmpdir:
            source = Path(tmpdir) / "logo.png"
            source.write_bytes(PNG_BYTES)
            state = SlideState("ppt/slides/slide1.xml")

            pic = emit_image(
                {
                    "name": "Logo",
                    "src": str(source),
                    "x": "1in",
                    "y": "1in",
                    "w": "1in",
                    "h": "1in",
                },
                state,
                relationships,
                content_types,
            )

        self.assertEqual(pic.tag, qn("p", "pic"))
        self.assertEqual(pic.find("p:nvPicPr/p:cNvPr", NSMAP).get("id"), "2")
        self.assertEqual(
            pic.find("p:blipFill/a:blip", NSMAP).get(qn("r", "embed")),
            "rId1",
        )
        self.assertEqual(state.media_parts["ppt/media/image1.png"], PNG_BYTES)
        self.assertEqual(content_types._defaults["png"], "image/png")

        rel = rels_for(relationships, "ppt/slides/slide1.xml")[0]
        self.assertEqual(rel.get("Type"), RelationshipRegistry.IMAGE)
        self.assertEqual(rel.get("Target"), "../media/image1.png")

    def test_image_placeholder_omits_xfrm_and_sets_placeholder_metadata(self):
        content_types = ContentTypeRegistry()
        relationships = RelationshipRegistry()

        with tempfile.TemporaryDirectory() as tmpdir:
            source = Path(tmpdir) / "photo.png"
            source.write_bytes(PNG_BYTES)
            state = SlideState("ppt/slides/slide1.xml")

            pic = emit_image(
                {
                    "name": "Inherited Picture",
                    "src": str(source),
                    "idx": 1,
                    "placeholder_type": "pic",
                },
                state,
                relationships,
                content_types,
            )

        ph = pic.find("p:nvPicPr/p:nvPr/p:ph", NSMAP)
        self.assertEqual(ph.get("type"), "pic")
        self.assertEqual(ph.get("idx"), "1")
        self.assertIsNone(pic.find("p:spPr/a:xfrm", NSMAP))

    def test_image_cover_crop_and_media_dedup_reuse_registered_part(self):
        content_types = ContentTypeRegistry()
        relationships = RelationshipRegistry()

        with tempfile.TemporaryDirectory() as tmpdir:
            source = Path(tmpdir) / "photo.png"
            source.write_bytes(PNG_BYTES)
            state = SlideState("ppt/slides/slide1.xml")

            first = emit_image(
                {
                    "name": "First",
                    "src": str(source),
                    "x": "1in",
                    "y": "1in",
                    "w": "2in",
                    "h": "2in",
                    "crop": {"l": 25000, "r": 25000},
                },
                state,
                relationships,
                content_types,
            )
            second = emit_image(
                {
                    "name": "Second",
                    "src": str(source),
                    "x": "4in",
                    "y": "1in",
                    "w": "2in",
                    "h": "2in",
                },
                state,
                relationships,
                content_types,
            )

        self.assertEqual(len(state.media_parts), 1)
        self.assertEqual(
            first.find("p:blipFill/a:blip", NSMAP).get(qn("r", "embed")),
            second.find("p:blipFill/a:blip", NSMAP).get(qn("r", "embed")),
        )
        src_rect = first.find("p:blipFill/a:srcRect", NSMAP)
        self.assertEqual(src_rect.get("l"), "25000")
        self.assertEqual(src_rect.get("r"), "25000")
        self.assertEqual(len(rels_for(relationships, "ppt/slides/slide1.xml")), 1)

    def test_slide_background_emits_solid_token_and_hex(self):
        token_bg = emit_slide_background(
            {"kind": "solid", "color": "accent5"},
            SlideState("ppt/slides/slide1.xml"),
            RelationshipRegistry(),
            ContentTypeRegistry(),
        )
        hex_bg = emit_slide_background(
            {"kind": "solid", "color": "#0F172A"},
            SlideState("ppt/slides/slide1.xml"),
            RelationshipRegistry(),
            ContentTypeRegistry(),
        )

        self.assertEqual(token_bg.tag, qn("p", "bg"))
        self.assertEqual(
            token_bg.find("p:bgPr/a:solidFill/a:schemeClr", NSMAP).get("val"),
            "accent5",
        )
        self.assertEqual(
            hex_bg.find("p:bgPr/a:solidFill/a:srgbClr", NSMAP).get("val"),
            "0F172A",
        )

    def test_slide_background_image_registers_media_relationship_and_content_type(self):
        content_types = ContentTypeRegistry()
        relationships = RelationshipRegistry()

        with tempfile.TemporaryDirectory() as tmpdir:
            source = Path(tmpdir) / "cover.png"
            source.write_bytes(PNG_BYTES)
            state = SlideState("ppt/slides/slide1.xml")

            bg = emit_slide_background(
                {"kind": "image", "src": str(source)},
                state,
                relationships,
                content_types,
            )

        self.assertEqual(
            bg.find("p:bgPr/a:blipFill/a:blip", NSMAP).get(qn("r", "embed")),
            "rId1",
        )
        self.assertEqual(state.media_parts["ppt/media/image1.png"], PNG_BYTES)
        self.assertEqual(content_types._defaults["png"], "image/png")

        rel = rels_for(relationships, "ppt/slides/slide1.xml")[0]
        self.assertEqual(rel.get("Type"), RelationshipRegistry.IMAGE)
        self.assertEqual(rel.get("Target"), "../media/image1.png")


class SlideBuilderTests(unittest.TestCase):
    def test_slide_builder_emits_required_tree_relationships_ids_and_z_order(self):
        content_types = ContentTypeRegistry()
        relationships = RelationshipRegistry()
        slide_data = {
            "index": 1,
            "layout": "titleBody",
            "name": "Content",
            "shapes": [
                {
                    "type": "placeholder_text",
                    "idx": 0,
                    "placeholder_type": "title",
                    "name": "Title Placeholder",
                    "text": "Title",
                },
                {
                    "type": "text_box",
                    "name": "Body Text Box",
                    "x": "1in",
                    "y": "2in",
                    "w": "3in",
                    "h": "1in",
                    "text": "Body",
                },
                {
                    "type": "autoshape",
                    "name": "Accent Rectangle",
                    "x": "5in",
                    "y": "2in",
                    "w": "2in",
                    "h": "1in",
                    "fill": "accent1",
                    "line": False,
                },
            ],
        }

        xml = SlideBuilder(content_types, relationships).build(slide_data)
        root = parse_xml(xml)

        self.assertEqual(root.tag, qn("p", "sld"))
        self.assertEqual([child.tag for child in root], [qn("p", "cSld"), qn("p", "clrMapOvr")])
        self.assertIsNotNone(root.find("p:clrMapOvr/a:masterClrMapping", NSMAP))
        self.assertEqual(root.find("p:cSld/p:spTree/p:nvGrpSpPr/p:cNvPr", NSMAP).get("id"), "1")

        ids = c_nv_pr_ids(root)
        self.assertEqual(len(ids), len(set(ids)))
        self.assertEqual(ids, ["1", "2", "3", "4"])
        self.assertEqual(
            z_order_names(root),
            ["Title Placeholder", "Body Text Box", "Accent Rectangle"],
        )
        placeholder = root.xpath(
            "p:cSld/p:spTree/p:sp[p:nvSpPr/p:nvPr/p:ph]",
            namespaces=NSMAP,
        )[0]
        self.assertIsNone(placeholder.find("p:spPr/a:xfrm", NSMAP))

        rel = rels_for(relationships, "ppt/slides/slide1.xml")[0]
        self.assertEqual(rel.get("Type"), RelationshipRegistry.SLIDE_LAYOUT)
        self.assertEqual(rel.get("Target"), "../slideLayouts/slideLayout2.xml")
        self.assertEqual(
            content_types.overrides["ppt/slides/slide1.xml"],
            ContentTypeRegistry.SLIDE,
        )

    def test_slide_builder_places_background_before_shape_tree(self):
        content_types = ContentTypeRegistry()
        relationships = RelationshipRegistry()
        slide_data = {
            "index": 1,
            "layout": "blank",
            "background": {"kind": "solid", "color": "accent5"},
            "shapes": [],
        }

        xml = SlideBuilder(content_types, relationships).build(slide_data)
        root = parse_xml(xml)
        c_sld = root.find("p:cSld", NSMAP)

        self.assertEqual([child.tag for child in c_sld], [qn("p", "bg"), qn("p", "spTree")])
        self.assertEqual(
            c_sld.find("p:bg/p:bgPr/a:solidFill/a:schemeClr", NSMAP).get("val"),
            "accent5",
        )


class ZIPAssemblerTests(unittest.TestCase):
    def test_assembler_writes_rels_shell_parts_and_content_types_last(self):
        graph = build_core_parts(synthetic_slide_count=0)

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "assembled.pptx"
            ZIPAssembler(
                graph["content_types"],
                graph["relationships"],
            ).assemble(output_path, graph["parts"])

            with ZipFile(output_path) as pptx:
                names = pptx.namelist()
                self.assertEqual(names[-1], "[Content_Types].xml")
                self.assertIn("_rels/.rels", names)
                self.assertIn("ppt/presProps.xml", names)
                self.assertIn("ppt/viewProps.xml", names)
                self.assertIn("ppt/tableStyles.xml", names)
                self.assertIn("docProps/core.xml", names)
                self.assertIn("docProps/app.xml", names)
                self.assertFalse(any(name.endswith("/") for name in names))
                for name in names:
                    if name.endswith(".xml"):
                        first_line = pptx.read(name).decode("utf-8").splitlines()[0]
                        self.assertEqual(first_line, CANONICAL_XML_DECLARATION)

                content_types = parse_xml(pptx.read("[Content_Types].xml").decode("utf-8"))
                for part_name in (
                    "/ppt/presProps.xml",
                    "/ppt/viewProps.xml",
                    "/ppt/tableStyles.xml",
                    "/docProps/core.xml",
                    "/docProps/app.xml",
                ):
                    self.assertIsNotNone(
                        content_types.find(f'ct:Override[@PartName="{part_name}"]', CT_NS)
                    )

                package_rels = parse_xml(pptx.read("_rels/.rels").decode("utf-8"))
                office_rels = package_rels.findall("rel:Relationship", REL_NS)
                self.assertEqual(len(office_rels), 3)
                targets_by_type = {rel.get("Type"): rel.get("Target") for rel in office_rels}
                self.assertEqual(
                    targets_by_type[RelationshipRegistry.OFFICE_DOC],
                    "ppt/presentation.xml",
                )
                self.assertEqual(
                    targets_by_type[RelationshipRegistry.CORE_PROPERTIES],
                    "docProps/core.xml",
                )
                self.assertEqual(
                    targets_by_type[RelationshipRegistry.EXTENDED_PROPERTIES],
                    "docProps/app.xml",
                )

                core_props = parse_xml(pptx.read("docProps/core.xml").decode("utf-8"))
                core_ns = {
                    "dc": "http://purl.org/dc/elements/1.1/",
                    "dcterms": "http://purl.org/dc/terms/",
                }
                self.assertEqual(
                    core_props.find("dc:title", core_ns).text,
                    "Paro Generated Presentation",
                )
                self.assertEqual(
                    core_props.find("dc:creator", core_ns).text,
                    "Paro PPTX Engine",
                )
                self.assertEqual(
                    core_props.find("dcterms:created", core_ns).text,
                    "2026-06-01T00:00:00Z",
                )
                self.assertEqual(
                    core_props.find("dcterms:modified", core_ns).text,
                    "2026-06-01T00:00:00Z",
                )

                presentation_rels = parse_xml(
                    pptx.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
                )
                rel_types = {
                    rel.get("Type")
                    for rel in presentation_rels.findall("rel:Relationship", REL_NS)
                }
                self.assertIn(RelationshipRegistry.PRES_PROPS, rel_types)
                self.assertIn(RelationshipRegistry.VIEW_PROPS, rel_types)
                self.assertIn(RelationshipRegistry.TABLE_STYLES, rel_types)

    def test_assembler_app_props_slide_count_matches_slide_parts(self):
        graph = build_core_parts(synthetic_slide_count=2)

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "assembled.pptx"
            ZIPAssembler(
                graph["content_types"],
                graph["relationships"],
            ).assemble(output_path, graph["parts"])

            with ZipFile(output_path) as pptx:
                app_props = parse_xml(pptx.read("docProps/app.xml").decode("utf-8"))

        app_ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
        self.assertEqual(app_props.find("ep:Application", app_ns).text, "Paro PPTX Engine")
        self.assertEqual(app_props.find("ep:Slides", app_ns).text, "2")


class EndToEndDeckTests(unittest.TestCase):
    def test_sample_deck_opens_and_preserves_placeholder_inheritance(self):
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            self.fail("python-pptx is required for slide-layer integration tests")

        with tempfile.TemporaryDirectory() as tmpdir:
            graph = build_deck(Path(tmpdir) / "sample_deck.pptx")

            presentation = PptxPresentation(str(graph["output_path"]))
            self.assertEqual(len(presentation.slides), 2)
            self.assertEqual(
                int(presentation.slide_width),
                REFERENCE["slide_sizes"]["16:9"]["cx"],
            )
            self.assertEqual(
                int(presentation.slide_height),
                REFERENCE["slide_sizes"]["16:9"]["cy"],
            )

            title_slide = presentation.slides[0]
            self.assertEqual(title_slide.shapes.title.text, "Quarterly Update")
            placeholder_texts = [shape.text for shape in title_slide.placeholders]
            self.assertIn("Generated by Paro", placeholder_texts)

            content_slide = presentation.slides[1]
            names = [shape.name for shape in content_slide.shapes]
            self.assertIn("Body Text Box", names)
            self.assertIn("Accent Rectangle", names)
            non_placeholder_names = [
                shape.name
                for shape in content_slide.shapes
                if not getattr(shape, "is_placeholder", False)
            ]
            self.assertIn("Body Text Box", non_placeholder_names)
            self.assertIn("Accent Rectangle", non_placeholder_names)

            title_root = parse_xml(graph["parts"]["ppt/slides/slide1.xml"])
            placeholders = title_root.findall("p:cSld/p:spTree/p:sp", NSMAP)
            for placeholder in placeholders:
                self.assertIsNotNone(placeholder.find("p:nvSpPr/p:nvPr/p:ph", NSMAP))
                self.assertIsNone(placeholder.find("p:spPr/a:xfrm", NSMAP))

            expected_orders = {
                "ppt/slides/slide1.xml": ["Title Placeholder", "Subtitle Placeholder"],
                "ppt/slides/slide2.xml": [
                    "Content Title Placeholder",
                    "Body Text Box",
                    "Accent Rectangle",
                ],
            }
            for slide_path, expected_order in expected_orders.items():
                with self.subTest(slide=slide_path):
                    root = parse_xml(graph["parts"][slide_path])
                    ids = c_nv_pr_ids(root)
                    self.assertEqual(len(ids), len(set(ids)))
                    self.assertEqual(z_order_names(root), expected_order)


if __name__ == "__main__":
    unittest.main()
