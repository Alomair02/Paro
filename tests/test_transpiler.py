import tempfile
import unittest
from pathlib import Path
from zipfile import ZipFile

from PIL import Image

from builders.slide_builder import SlideBuilder
from core.content_type_reg import ContentTypeRegistry
from core.relationship_reg import RelationshipRegistry
from core.xml_builder import NSMAP, qn
from builders.common import REFERENCE
from transpiler import (
    DSLParser,
    LayoutRegistry,
    LayoutResolver,
    ShapeLibrary,
    ThemeRegistry,
    TimelineStyleRegistry,
    Validator,
    transpile_deck,
)
from transpiler.registries import ChartStyleRegistry
from transpiler.validator import TranspileValidationError
from transpiler.pipeline import _inline_theme_map
from tests.pptx_test_utils import parse_xml
from utils.converter import UnitConverter


def parse_resolve(xml: str):
    ast = DSLParser().parse(xml)
    registry = ThemeRegistry(
        {
            ast.inline_theme.name: {
                "name": ast.inline_theme.name,
                "colors": ThemeRegistry().get("default")["colors"],
                "fonts": ThemeRegistry().get("default")["fonts"],
            }
        }
        if ast.inline_theme
        else None
    )
    return LayoutResolver(registry, LayoutRegistry()).resolve(ast)


def shape_by_name(slide_data: dict, name: str):
    return next(shape for shape in slide_data["shapes"] if shape.get("name") == name)

def parse_xml_localname(element) -> str:
    from lxml import etree
    return etree.QName(element).localname

def write_test_image(path: Path, size: tuple[int, int] = (200, 100)):
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, (37, 99, 235)).save(path)


class TranspilerRegistryTests(unittest.TestCase):
    def test_theme_and_layout_registries_resolve_through_get(self):
        theme = ThemeRegistry().get("default")
        layout = LayoutRegistry().get("titleBody")

        self.assertEqual(theme["name"], "default")
        self.assertEqual(set(theme["typeScale"]), set(REFERENCE["default_theme"]["typeScale"]))
        self.assertEqual(layout["part_path"], "ppt/slideLayouts/slideLayout2.xml")

        with self.assertRaises(KeyError):
            ThemeRegistry().get("missing")
        with self.assertRaises(KeyError):
            LayoutRegistry().get("missing")

    def test_theme_registry_type_scale_can_be_overridden_per_role(self):
        custom = ThemeRegistry().get("default")
        custom["typeScale"]["body"] = {
            "size": 19,
            "weight": "bold",
            "spaceAfter": "8pt",
            "lineSpacing": 1.2,
        }

        theme = ThemeRegistry({"custom": custom}).get("custom")

        self.assertEqual(theme["typeScale"]["body"]["size"], 19)
        self.assertEqual(theme["typeScale"]["body"]["weight"], "bold")
        for role in ("title", "heading", "subheading", "body", "bodySmall", "caption"):
            self.assertIn("size", theme["typeScale"][role])
    
    def test_registries_merge_custom_over_defaults_keeping_builtins(self):
        # Uniform rule: passing custom data ADDS to / overrides built-ins;
        # built-ins are never lost. Same semantics across layout/chart/timeline.

        # LayoutRegistry: built-in survives, custom is added
        lr = LayoutRegistry({"my_layout": {"placeholders": []}})
        self.assertIsNotNone(lr.get("my_layout"))           # custom present
        self.assertIsNotNone(lr.get("titleBody"))           # a built-in still present
                                                            # (use any real default name)

        # ChartStyleRegistry: built-in survives, custom added
        cr = ChartStyleRegistry(styles={"my_style": {"palette": ["accent1"]}})
        self.assertIsNotNone(cr.get("my_style"))
        self.assertIsNotNone(cr.get("clean"))               # DEFAULT_STYLE survives

        # TimelineStyleRegistry: built-in survives, custom added
        tr = TimelineStyleRegistry(styles={"my_tl": {"mechanism": "bars-in-columns"}})
        self.assertIsNotNone(tr.get("my_tl"))
        self.assertIsNotNone(tr.get("gantt-grid"))          # DEFAULT_STYLE survives

    def test_custom_entry_overrides_builtin_by_name(self):
        # Reusing a built-in name overrides it (the deliberate-override path)
        cr = ChartStyleRegistry(styles={"clean": {"palette": ["accent6"]}})
        self.assertEqual(cr.get("clean")["palette"], ["accent6"])

    def test_timeline_style_registry_resolves_initial_style_bundles(self):
        registry = TimelineStyleRegistry()

        expected = {
            "gantt-grid": ("bars-in-columns", "gutter", "extend"),
            "gantt-minimal": ("bars-in-columns", "gutter", "stagger"),
            "roadmap": ("bars-in-columns", "on-bar", "extend"),
            "swimlane": ("bars-in-columns", "gutter", "stagger"),
        }
        for name, (mechanism, label_placement, milestone_overflow) in expected.items():
            style = registry.get(name)
            self.assertEqual(style["mechanism"], mechanism)
            self.assertEqual(style["labelPlacement"], label_placement)
            self.assertEqual(style["milestoneLabelOverflow"], milestone_overflow)
            self.assertIn("defaultTones", style)
            self.assertIn("finish", style)

        self.assertEqual(registry.get()["mechanism"], "bars-in-columns")
        for finish in ("flat", "soft-shadow", "elevated", "outlined", "accent-gradient"):
            finish_def = registry.get_finish(finish)
            self.assertIn("fill", finish_def)
            self.assertIn("cornerRadius", finish_def)
        with self.assertRaises(KeyError):
            registry.get("missing")
        with self.assertRaises(KeyError):
            registry.get_finish("missing")


class DSLParserTests(unittest.TestCase):
    def test_parser_round_trips_supported_elements_and_expands_defs(self):
        xml = """
        <deck theme="default" size="16:9" font="Aptos" background="accent5">
          <theme name="custom" accent1="#0057B8" heading="Arial" body="Arial"/>
          <defs>
            <def name="badge"><shape fill="accent1"><text>Badge</text></shape></def>
            <def name="footer" auto="true"><text size="8pt">Footer</text></def>
          </defs>
          <slide layout="blank" flow="free" background="#0F172A">
            <use ref="badge"/>
            <stack x="0in" y="0in" w="3in" h="1in"><text role="heading"><p role="caption"><run bold="true">Run</run> tail</p></text></stack>
            <grid x="0in" y="1in" w="3in" h="1in" cols="2"><shape col="2" fill="accent2"/></grid>
            <free x="0in" y="2in" w="4in" h="2in">
              <line x1="0%" y1="0%" x2="100%" y2="0%" color="dk1"/>
              <image src="logo.png" x="0in" y="0.1in" w="0.5in" h="0.5in"/>
              <table cols="2"><row><cell bold="true">A</cell><cell>B</cell></row></table>
              <timeline periods="2" finish="outlined" borderWidth="1.5pt" shadow="false"><group label="Delivery"><task label="Build" start="1"/></group><milestone label="Go" at="2"/></timeline>
            </free>
          </slide>
        </deck>
        """

        deck = DSLParser().parse(xml)

        self.assertEqual(deck.theme, "default")
        self.assertEqual(deck.background, "accent5")
        self.assertEqual(deck.inline_theme.name, "custom")
        self.assertEqual(deck.slides[0].layout, "blank")
        self.assertEqual(deck.slides[0].attrs["background"], "#0F172A")
        kinds = [child.kind for child in deck.slides[0].children]
        self.assertEqual(kinds[0], "shape")
        self.assertEqual(kinds[-1], "text")
        nested = []

        def walk(node):
            nested.append(node.kind)
            for child in node.children:
                walk(child)

        for child in deck.slides[0].children:
            walk(child)
        for expected in ("stack", "grid", "free", "line", "image", "table", "row", "cell", "timeline", "group", "task", "milestone", "run"):
            self.assertIn(expected, nested)
        stack = next(child for child in deck.slides[0].children if child.kind == "stack")
        self.assertEqual(stack.children[0].attrs["role"], "heading")
        self.assertEqual(stack.children[0].children[0].attrs["role"], "caption")

    def test_unknown_element_or_attribute_is_rejected(self):
        with self.assertRaisesRegex(Exception, "Unknown element"):
            DSLParser().parse("<deck><slide><bogus/></slide></deck>")
        with self.assertRaisesRegex(Exception, "Unknown attribute"):
            DSLParser().parse('<deck madeup="1"><slide/></deck>')
        with self.assertRaisesRegex(Exception, "Unknown type role"):
            DSLParser().parse('<deck><slide><text role="tiny">Bad</text></slide></deck>')


class LayoutResolverTests(unittest.TestCase):
    def test_role_resolution_uses_theme_type_scale_and_paragraph_override(self):
        custom = ThemeRegistry().get("default")
        custom["typeScale"]["heading"] = {
            "size": 23,
            "weight": "bold",
            "spaceAfter": "7pt",
            "lineSpacing": 1.2,
        }
        custom["typeScale"]["caption"] = {
            "size": 12,
            "weight": "normal",
            "spaceAfter": "2pt",
            "lineSpacing": 1.05,
        }
        ast = DSLParser().parse(
            """
            <deck theme="custom">
              <slide layout="blank" flow="free">
                <text x="1in" y="1in" w="4in" h="1.5in" role="heading">
                  <p>Heading</p>
                  <p role="caption">Caption</p>
                </text>
              </slide>
            </deck>
            """
        )

        deck = LayoutResolver(ThemeRegistry({"custom": custom}), LayoutRegistry()).resolve(ast)

        paragraphs = deck.slide_data[0]["shapes"][0]["paragraphs"]
        self.assertEqual(paragraphs[0]["runs"][0]["size_pt"], 23)
        self.assertTrue(paragraphs[0]["runs"][0]["bold"])
        self.assertEqual(paragraphs[0]["spaceAfter"], "7pt")
        self.assertEqual(paragraphs[0]["lineSpacing"], 1.2)
        self.assertEqual(paragraphs[1]["runs"][0]["size_pt"], 12)
        self.assertNotIn("bold", paragraphs[1]["runs"][0])

    def test_role_size_omits_when_placeholder_master_already_matches(self):
        custom = ThemeRegistry().get("default")
        custom["typeScale"]["title"]["size"] = REFERENCE["master_text_styles"]["title"]["size"]
        ast = DSLParser().parse(
            """
            <deck theme="custom">
              <slide layout="title">
                <text placeholder="title" role="title">Inherited title size</text>
              </slide>
            </deck>
            """
        )

        deck = LayoutResolver(ThemeRegistry({"custom": custom}), LayoutRegistry()).resolve(ast)

        run = deck.slide_data[0]["shapes"][0]["paragraphs"][0]["runs"][0]
        self.assertNotIn("size_pt", run)

    def test_stack_edge_anchor_computes_full_bleed_thirds(self):
        deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <stack dir="h" anchor="bottom" gap="0" h="0.3in">
                  <shape fill="#E53935"/>
                  <shape fill="#FB8C00"/>
                  <shape fill="#43A047"/>
                </stack>
              </slide>
            </deck>
            """
        )

        shapes = deck.slide_data[0]["shapes"]
        slide_w = REFERENCE["slide_sizes"]["16:9"]["cx"]
        expected_h = UnitConverter.to_emu("0.3in")
        self.assertEqual([shape["x"] for shape in shapes], [0, round(slide_w / 3), round(slide_w * 2 / 3)])
        self.assertEqual(shapes[0]["y"], REFERENCE["slide_sizes"]["16:9"]["cy"] - expected_h)
        self.assertEqual(shapes[0]["w"], round(slide_w / 3))

    def test_grid_resolver_computes_cell_boxes(self):
        deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <grid x="1in" y="1in" w="4in" h="2in" cols="2" rows="1" gap="0in">
                  <shape col="2" fill="accent1"/>
                </grid>
              </slide>
            </deck>
            """
        )

        shape = deck.slide_data[0]["shapes"][0]
        self.assertEqual(shape["x"], UnitConverter.to_emu("3in"))
        self.assertEqual(shape["y"], UnitConverter.to_emu("1in"))
        self.assertEqual(shape["w"], UnitConverter.to_emu("2in"))
        self.assertEqual(shape["h"], UnitConverter.to_emu("2in"))

    def test_table_resolver_emits_native_table_shape_data(self):
        deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <table x="1in" y="1in" w="4in" h="1in" cols="2" header="true" colWidths="50%,50%">
                  <row><cell>A</cell><cell>B</cell></row>
                </table>
              </slide>
            </deck>
            """
        )

        shapes = deck.slide_data[0]["shapes"]
        self.assertEqual([shape["type"] for shape in shapes], ["table"])
        self.assertEqual(shapes[0]["columns"][0]["w"], UnitConverter.to_emu("2in"))
        self.assertEqual(shapes[0]["rows"][0]["cells"][0]["fill"], "accent1")

    def test_timeline_resolver_aligns_start_and_span_to_period_tracks(self):
        deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="6in" h="2in" periods="6" style="gantt-grid">
                  <task label="Build" start="3" span="4" tone="accent2"/>
                </timeline>
              </slide>
            </deck>
            """
        )

        bar = shape_by_name(deck.slide_data[0], "Timeline Task Build")
        style = TimelineStyleRegistry().get("gantt-grid")
        actual_x = UnitConverter.to_emu("1in")
        actual_w = UnitConverter.to_emu("6in")
        gutter = max(round(actual_w * style["gutterRatio"]), UnitConverter.to_emu(style["gutterMin"]))
        plot_x = actual_x + gutter
        period_w = (actual_w - gutter) / 6
        inset = UnitConverter.to_emu(style["barInset"])
        expected_left = round(plot_x + 2 * period_w) + inset
        expected_right = round(plot_x + 6 * period_w) - inset
        self.assertEqual(bar["x"], expected_left)
        self.assertEqual(bar["w"], expected_right - expected_left)

        timeline_text_shapes = [
            shape
            for shape in deck.slide_data[0]["shapes"]
            if shape.get("name") in {"Timeline Period", "Timeline Task Label"}
        ]
        self.assertTrue(timeline_text_shapes)
        for shape in timeline_text_shapes:
            self.assertEqual(shape.get("wrap"), "none")
            for paragraph in shape["paragraphs"]:
                for run in paragraph.get("runs", []):
                    self.assertIn("size_pt", run)

    def test_timeline_geometry_is_compact_and_milestones_are_below_header(self):
        deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="7in" h="4in" periods="4" style="gantt-grid">
                  <task label="One" start="1" span="1"/>
                  <task label="Two" start="2" span="1"/>
                  <task label="Three" start="3" span="1"/>
                  <milestone label="Review" at="2"/>
                </timeline>
              </slide>
            </deck>
            """
        )

        style = TimelineStyleRegistry().get("gantt-grid")
        row_h = UnitConverter.to_emu(style["rowHeight"])
        label_shapes = [
            shape for shape in deck.slide_data[0]["shapes"] if shape.get("name") == "Timeline Task Label"
        ]
        bars = [
            shape
            for shape in deck.slide_data[0]["shapes"]
            if shape.get("name", "").startswith("Timeline Task ") and shape.get("name") != "Timeline Task Label"
        ]
        milestones = [
            shape for shape in deck.slide_data[0]["shapes"] if shape.get("name", "").startswith("Timeline Milestone ")
        ]
        header_bottom = UnitConverter.to_emu("1in") + UnitConverter.to_emu(style["headerHeight"])
        task_band_bottom = header_bottom + 3 * row_h

        self.assertEqual([shape["h"] for shape in label_shapes], [row_h, row_h, row_h])
        self.assertLessEqual(max(shape["y"] + shape["h"] for shape in bars), task_band_bottom)
        self.assertGreaterEqual(min(shape["y"] for shape in milestones), task_band_bottom)
        self.assertGreaterEqual(min(shape["y"] for shape in milestones), header_bottom)

    def test_timeline_group_data_model_and_unit_labels_parse_and_resolve(self):
        ast = DSLParser().parse(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="8in" h="3in" periods="2" unit="month" style="swimlane">
                  <group label="Discovery">
                    <task label="Research" start="1" span="1" tone="accent2"/>
                  </group>
                  <group label="Delivery">
                    <task label="Build" start="2" span="1"/>
                  </group>
                  <milestone label="Review" at="2"/>
                </timeline>
              </slide>
            </deck>
            """
        )

        timeline = ast.slides[0].children[0]
        self.assertEqual(timeline.attrs["unit"], "month")
        self.assertEqual(timeline.children[0].kind, "group")
        self.assertEqual(timeline.children[0].children[0].attrs["tone"], "accent2")

        deck = LayoutResolver().resolve(ast)
        self.assertTrue(shape_by_name(deck.slide_data[0], "Timeline Group Discovery"))
        period_labels = [
            shape["paragraphs"][0]["runs"][0]["text"]
            for shape in deck.slide_data[0]["shapes"]
            if shape.get("name") == "Timeline Period"
        ]
        self.assertEqual(period_labels, ["Month 1", "Month 2"])

    def test_timeline_label_truncate_and_shrink_never_wrap(self):
        base = TimelineStyleRegistry().get("gantt-grid")
        shrink = dict(base)
        shrink["labelOverflow"] = "shrink"
        registry = TimelineStyleRegistry({"compact-shrink": shrink})
        ast = DSLParser().parse(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="2.2in" h="1.4in" periods="2" style="compact-shrink">
                  <task label="Extremely long task label that would normally wrap" start="1" span="2"/>
                </timeline>
              </slide>
            </deck>
            """
        )
        deck = LayoutResolver(timeline_style_registry=registry).resolve(ast)

        label = next(shape for shape in deck.slide_data[0]["shapes"] if shape.get("name") == "Timeline Task Label")
        paragraph = label["paragraphs"][0]
        run = paragraph["runs"][0]
        measurement = Validator().measurer.measure(
            run["text"],
            ThemeRegistry().get("default")["fonts"]["body"],
            run["size_pt"],
            label["w"],
        )
        self.assertLessEqual(measurement.wrapped_lines, 1)
        self.assertLessEqual(measurement.rendered_width_emu, label["w"])
        self.assertLessEqual(run["size_pt"], REFERENCE["default_theme"]["typeScale"]["bodySmall"]["size"])
        self.assertGreaterEqual(run["size_pt"], REFERENCE["default_theme"]["typeScale"]["bodySmall"]["minSize"])

    def test_timeline_milestone_extend_and_stagger_keep_labels_clear(self):
        extend_deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="5in" h="2in" periods="5" style="gantt-grid">
                  <task label="Build" start="1" span="5"/>
                  <milestone label="Executive steering committee" at="3"/>
                </timeline>
              </slide>
            </deck>
            """
        )
        extend_label = next(shape for shape in extend_deck.slide_data[0]["shapes"] if shape.get("name") == "Timeline Milestone Label")
        style = TimelineStyleRegistry().get("gantt-grid")
        timeline_w = UnitConverter.to_emu("5in")
        gutter = max(round(timeline_w * style["gutterRatio"]), UnitConverter.to_emu(style["gutterMin"]))
        period_w = round((timeline_w - gutter) / 5)
        self.assertGreater(extend_label["w"], period_w)

        stagger_deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="5in" h="2in" periods="5" style="gantt-minimal">
                  <task label="Build" start="1" span="5"/>
                  <milestone label="Alpha" at="3"/>
                  <milestone label="Beta" at="3"/>
                </timeline>
              </slide>
            </deck>
            """
        )
        labels = [shape for shape in stagger_deck.slide_data[0]["shapes"] if shape.get("name") == "Timeline Milestone Label"]
        self.assertEqual(len(labels), 2)
        self.assertNotEqual(labels[0]["y"], labels[1]["y"])
        self.assertFalse(
            labels[0]["x"] < labels[1]["x"] + labels[1]["w"]
            and labels[1]["x"] < labels[0]["x"] + labels[0]["w"]
            and labels[0]["y"] < labels[1]["y"] + labels[1]["h"]
            and labels[1]["y"] < labels[0]["y"] + labels[0]["h"]
        )

    def test_timeline_finishes_emit_distinct_sppr_and_overrides(self):
        ast = DSLParser().parse(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="0.5in" y="0.5in" w="3in" h="1.2in" periods="2" style="gantt-grid" finish="flat">
                  <task label="flat" start="1" span="1"/>
                </timeline>
                <timeline x="4in" y="0.5in" w="3in" h="1.2in" periods="2" style="gantt-grid" finish="soft-shadow">
                  <task label="soft" start="1" span="1"/>
                </timeline>
                <timeline x="0.5in" y="2in" w="3in" h="1.2in" periods="2" style="gantt-grid" finish="elevated">
                  <task label="elevated" start="1" span="1"/>
                </timeline>
                <timeline x="4in" y="2in" w="3in" h="1.2in" periods="2" style="gantt-grid" finish="outlined">
                  <task label="outlined" start="1" span="1"/>
                </timeline>
                <timeline x="0.5in" y="3.5in" w="3in" h="1.2in" periods="2" style="gantt-grid" finish="accent-gradient">
                  <task label="gradient" start="1" span="1"/>
                </timeline>
                <timeline x="4in" y="3.5in" w="3in" h="1.2in" periods="2" style="roadmap" finish="flat" borderWidth="2pt" shadow="true">
                  <task label="override" start="1" span="1"/>
                </timeline>
              </slide>
            </deck>
            """
        )
        deck = LayoutResolver().resolve(ast)
        slide_xml = SlideBuilder(
            ContentTypeRegistry(),
            RelationshipRegistry(),
            {},
        ).build(deck.slide_data[0], part_path="ppt/slides/slide1.xml")
        root = parse_xml(slide_xml)

        def sppr_for(task_name: str):
            for sp in root.findall(".//p:sp", NSMAP):
                c_nv_pr = sp.find("p:nvSpPr/p:cNvPr", NSMAP)
                if c_nv_pr is not None and c_nv_pr.get("name") == f"Timeline Task {task_name}":
                    return sp.find("p:spPr", NSMAP)
            raise AssertionError(f"missing task {task_name}")

        flat = sppr_for("flat")
        soft = sppr_for("soft")
        elevated = sppr_for("elevated")
        outlined = sppr_for("outlined")
        gradient = sppr_for("gradient")
        override = sppr_for("override")

        self.assertIsNone(flat.find("a:effectLst/a:outerShdw", NSMAP))
        self.assertIsNotNone(soft.find("a:effectLst/a:outerShdw", NSMAP))
        self.assertIsNotNone(elevated.find("a:ln/a:solidFill", NSMAP))
        self.assertEqual(elevated.find("a:ln", NSMAP).get("w"), str(UnitConverter.to_emu("0.75pt")))
        self.assertIsNotNone(outlined.find("a:solidFill/a:schemeClr/a:lumOff", NSMAP))
        self.assertIsNotNone(gradient.find("a:gradFill", NSMAP))
        self.assertEqual(override.find("a:ln", NSMAP).get("w"), str(UnitConverter.to_emu("2pt")))
        self.assertIsNotNone(override.find("a:effectLst/a:outerShdw", NSMAP))
        self.assertIsNone(override.find("a:gradFill", NSMAP))

        gd_values = {
            task: sppr_for(task).find("a:prstGeom/a:avLst/a:gd", NSMAP).get("fmla")
            for task in ("flat", "soft", "gradient")
        }
        self.assertEqual(len(set(gd_values.values())), 3)

    def test_timeline_swimlane_renders_group_bands_and_falls_back_without_groups(self):
        grouped = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="7in" h="3in" periods="3" style="swimlane">
                  <group label="Discovery"><task label="Research" start="1" span="1"/></group>
                  <group label="Delivery"><task label="Build" start="2" span="2"/></group>
                </timeline>
              </slide>
            </deck>
            """
        )
        self.assertTrue(shape_by_name(grouped.slide_data[0], "Timeline Group Discovery"))
        self.assertTrue(shape_by_name(grouped.slide_data[0], "Timeline Group Delivery"))

        flat = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <timeline x="1in" y="1in" w="7in" h="3in" periods="3" style="swimlane">
                  <task label="Research" start="1" span="1"/>
                </timeline>
              </slide>
            </deck>
            """
        )
        self.assertFalse([shape for shape in flat.slide_data[0]["shapes"] if shape.get("name", "").startswith("Timeline Group")])
        self.assertTrue(shape_by_name(flat.slide_data[0], "Timeline Task Research"))

    def test_alignment_round_trips_into_slide_xml(self):
        deck = parse_resolve(
            """
            <deck>
              <slide layout="blank" flow="free">
                <text x="0.5in" y="0.5in" w="2in" h="0.4in" align="l">Left</text>
                <text x="0.5in" y="1.0in" w="2in" h="0.4in" align="ctr">Center</text>
                <text x="0.5in" y="1.5in" w="2in" h="0.4in" align="r">Right</text>
                <text x="0.5in" y="2.0in" w="2in" h="0.4in" align="just">Justified</text>
              </slide>
            </deck>
            """
        )
        slide_xml = SlideBuilder(
            ContentTypeRegistry(),
            RelationshipRegistry(),
            {},
        ).build(deck.slide_data[0], part_path="ppt/slides/slide1.xml")

        for align in ("l", "ctr", "r", "just"):
            self.assertIn(f'algn="{align}"', slide_xml)

    def test_named_image_assets_resolve_through_shape_library(self):
        ast = DSLParser().parse(
            """
            <deck>
              <slide layout="blank" flow="free">
                <image src="brand_logo" x="1in" y="1in" w="1in" h="1in"/>
              </slide>
            </deck>
            """
        )
        deck = LayoutResolver(
            ThemeRegistry(),
            LayoutRegistry(),
            ShapeLibrary({"brand_logo": {"src": "assets/logo.png", "alt": "Brand Logo"}}),
        ).resolve(ast)

        image = deck.slide_data[0]["shapes"][0]
        self.assertEqual(image["src"], "assets/logo.png")
        self.assertEqual(image["name"], "Brand Logo")

    def test_image_resolver_computes_managed_free_placeholder_and_fit_modes(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            image_path = Path(tmpdir) / "wide.png"
            write_test_image(image_path, (200, 100))
            deck = parse_resolve(
                f"""
                <deck>
                  <slide layout="picture" flow="free">
                    <image placeholder="pic" src="{image_path}" fit="cover"/>
                    <image src="{image_path}" x="1in" y="1in" w="2in" h="2in" fit="cover"/>
                    <image src="{image_path}" x="4in" y="1in" w="2in" h="2in" fit="stretch"/>
                    <grid x="1in" y="4in" w="4in" h="2in" cols="2" rows="1" gap="0in">
                      <image src="{image_path}" col="2" fit="contain"/>
                    </grid>
                  </slide>
                </deck>
                """
            )

        placeholder, cover, stretch, managed = deck.slide_data[0]["shapes"]

        self.assertEqual(placeholder["placeholder_type"], "pic")
        self.assertEqual(placeholder["idx"], 1)
        self.assertNotIn("x", placeholder)
        self.assertEqual(placeholder["crop"], {"l": 0, "t": 8620, "r": 0, "b": 8620})

        self.assertEqual(cover["x"], UnitConverter.to_emu("1in"))
        self.assertEqual(cover["w"], UnitConverter.to_emu("2in"))
        self.assertEqual(cover["crop"], {"l": 25000, "t": 0, "r": 25000, "b": 0})

        self.assertEqual(stretch["x"], UnitConverter.to_emu("4in"))
        self.assertNotIn("crop", stretch)

        self.assertEqual(managed["x"], UnitConverter.to_emu("3in"))
        self.assertEqual(managed["y"], UnitConverter.to_emu("4.5in"))
        self.assertEqual(managed["w"], UnitConverter.to_emu("2in"))
        self.assertEqual(managed["h"], UnitConverter.to_emu("1in"))

    def test_background_inherits_from_deck_and_slide_overrides(self):
        deck = parse_resolve(
            """
            <deck background="accent5">
              <slide layout="blank" flow="free"/>
              <slide layout="blank" flow="free" background="#0F172A"/>
            </deck>
            """
        )
        no_background = parse_resolve('<deck><slide layout="blank" flow="free"/></deck>')

        self.assertEqual(
            deck.slide_data[0]["background"],
            {"kind": "solid", "color": "accent5"},
        )
        self.assertEqual(
            deck.slide_data[1]["background"],
            {"kind": "solid", "color": "#0F172A"},
        )
        self.assertNotIn("background", no_background.slide_data[0])


class TranspilerValidatorTests(unittest.TestCase):
    def test_validator_reports_bounds_grid_placeholder_overlap_and_font_drift(self):
        bounds = parse_resolve(
            '<deck><slide layout="blank" flow="free"><text x="13in" y="0in" w="1in" h="1in">Too far</text></slide></deck>'
        )
        grid = parse_resolve(
            '<deck><slide layout="blank" flow="free"><grid x="0in" y="0in" w="2in" h="1in" cols="2" rows="1"><shape col="3" fill="accent1"/></grid></slide></deck>'
        )
        placeholder = parse_resolve(
            '<deck><slide layout="title" flow="stack"><text placeholder="body">Missing</text></slide></deck>'
        )
        overlap = parse_resolve(
            """
            <deck><slide layout="blank" flow="free">
              <shape x="1in" y="1in" w="2in" h="1in" fill="accent1"/>
              <shape x="1.5in" y="1in" w="2in" h="1in" fill="accent2"/>
            </slide></deck>
            """
        )
        drift = parse_resolve(
            """
            <deck><slide layout="blank" flow="stack">
              <text size="10pt">Small</text>
              <text size="18pt">Large</text>
            </slide></deck>
            """
        )

        validator = Validator()
        self.assertIn("bounds", {issue.code for issue in validator.validate(bounds)})
        self.assertIn("grid_bounds", {issue.code for issue in validator.validate(grid)})
        self.assertIn("placeholder", {issue.code for issue in validator.validate(placeholder)})
        self.assertIn("free_overlap", {issue.code for issue in validator.validate(overlap)})
        self.assertIn("font_drift", {issue.code for issue in validator.validate(drift)})
        with self.assertRaises(TranspileValidationError):
            validator.raise_for_errors(validator.validate(bounds))

    def test_validator_reports_missing_image_without_resolver_crash(self):
        missing = Path("tests/fixtures/does-not-exist.png")
        deck = parse_resolve(
            f"""
            <deck>
              <slide layout="blank" flow="free">
                <image src="{missing}" x="1in" y="1in" w="2in" h="1in"/>
              </slide>
            </deck>
            """
        )

        issues = Validator().validate(deck)

        missing_issues = [issue for issue in issues if issue.code == "missing_image"]
        self.assertEqual(len(missing_issues), 1)
        self.assertIn(str(missing), missing_issues[0].message)
        with self.assertRaises(TranspileValidationError):
            Validator().raise_for_errors(issues)


class TranspilerIntegrationTests(unittest.TestCase):
    def test_background_demo_transpiles_opens_and_emits_bg_first(self):
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            self.fail("python-pptx is required for transpiler integration tests")

        source_path = Path("tests/fixtures/golden/backgrounds_demo.xml")
        output_path = Path("tests/fixtures/backgrounds_demo.pptx")
        result = transpile_deck(source_path, output_path)

        presentation = PptxPresentation(str(result.pptx_path))
        self.assertEqual(len(presentation.slides), 2)
        self.assertFalse([issue for issue in result.validation_issues if issue.severity == "error"])

        with ZipFile(output_path) as pptx:
            slide1 = parse_xml(pptx.read("ppt/slides/slide1.xml").decode("utf-8"))
            slide2 = parse_xml(pptx.read("ppt/slides/slide2.xml").decode("utf-8"))

        for root in (slide1, slide2):
            c_sld = root.find("p:cSld", NSMAP)
            self.assertEqual([child.tag for child in c_sld[:2]], [qn("p", "bg"), qn("p", "spTree")])

        self.assertEqual(
            slide1.find("p:cSld/p:bg/p:bgPr/a:solidFill/a:schemeClr", NSMAP).get("val"),
            "accent5",
        )
        self.assertEqual(
            slide2.find("p:cSld/p:bg/p:bgPr/a:solidFill/a:srgbClr", NSMAP).get("val"),
            "0F172A",
        )

    def test_images_demo_transpiles_opens_and_contains_picture_parts(self):
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            self.fail("python-pptx is required for transpiler integration tests")

        source_path = Path("tests/fixtures/golden/images_demo.xml")
        output_path = Path("tests/fixtures/images_demo.pptx")
        result = transpile_deck(source_path, output_path)

        presentation = PptxPresentation(str(result.pptx_path))
        self.assertEqual(len(presentation.slides), 1)
        self.assertFalse([issue for issue in result.validation_issues if issue.severity == "error"])

        with ZipFile(output_path) as pptx:
            names = pptx.namelist()
            media_names = [name for name in names if name.startswith("ppt/media/")]
            slide = parse_xml(pptx.read("ppt/slides/slide1.xml").decode("utf-8"))

        self.assertEqual(media_names, ["ppt/media/image1.png"])
        pics = slide.findall("p:cSld/p:spTree/p:pic", NSMAP)
        self.assertEqual(len(pics), 3)
        self.assertIsNotNone(pics[0].find("p:nvPicPr/p:nvPr/p:ph", NSMAP))
        self.assertIsNone(pics[0].find("p:spPr/a:xfrm", NSMAP))
        self.assertIsNotNone(pics[1].find("p:blipFill/a:srcRect", NSMAP))

    def test_timeline_styles_demo_transpiles_opens_and_contains_timeline_primitives(self):
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            self.fail("python-pptx is required for transpiler integration tests")

        source_path = Path("tests/fixtures/golden/timeline_styles.xml")
        output_path = Path("tests/fixtures/timeline_styles.pptx")
        result = transpile_deck(source_path, output_path)

        presentation = PptxPresentation(str(result.pptx_path))
        self.assertEqual(len(presentation.slides), 6)
        self.assertFalse([issue for issue in result.validation_issues if issue.severity == "error"])

        with ZipFile(output_path) as pptx:
            slide_xml = "\n".join(
                pptx.read(name).decode("utf-8")
                for name in pptx.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )
            luminance_values = []
            connector_extents = []
            for name in pptx.namelist():
                if not name.endswith(".xml"):
                    continue
                root = parse_xml(pptx.read(name).decode("utf-8"))
                for transform in ("lumMod", "lumOff"):
                    luminance_values.extend(
                        int(node.get("val"))
                        for node in root.findall(f".//a:{transform}", NSMAP)
                    )
                if name.startswith("ppt/slides/slide"):
                    connector_extents.extend(
                        (int(ext.get("cx")), int(ext.get("cy")))
                        for ext in root.findall(".//p:cxnSp/p:spPr/a:xfrm/a:ext", NSMAP)
                    )

        self.assertIn("Timeline Task Discovery research", slide_xml)
        self.assertIn("Timeline Gridline", slide_xml)
        self.assertIn("Timeline Milestone Label", slide_xml)
        self.assertIn("Timeline Group Discovery", slide_xml)
        self.assertIn("Timeline Task Weekly Leads Meeting", slide_xml)
        self.assertIn("Project status meeting", slide_xml)
        self.assertIn("Timeline Task accent-gradient", slide_xml)
        self.assertIn("<a:gradFill", slide_xml)
        self.assertIn("<a:outerShdw", slide_xml)
        self.assertIn("<a:gd name=\"adj\"", slide_xml)
        self.assertTrue(luminance_values)
        self.assertLessEqual(max(luminance_values), 100000)
        self.assertTrue(connector_extents)
        self.assertTrue(all(cx > 0 and cy > 0 for cx, cy in connector_extents))
        self.assertIn("<p:cxnSp>", slide_xml)
        self.assertTrue(output_path.exists())

    def test_pressure_test_deck_transpiles_and_opens_with_python_pptx(self):
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            self.fail("python-pptx is required for transpiler integration tests")

        dsl = """
        <deck theme="pressure" size="16:9">
          <theme name="pressure"
                 dk1="#111827" lt1="#FFFFFF" dk2="#374151" lt2="#F3F4F6"
                 accent1="#2563EB" accent2="#F97316" accent3="#16A34A"
                 accent4="#7C3AED" accent5="#0E7490" accent6="#DC2626"
                 hlink="#1D4ED8" folHlink="#6D28D9"
                 heading="Aptos Display" body="Aptos"/>
          <slide layout="titleBody" flow="free">
            <text placeholder="title">Styled panel grid</text>
            <grid x="0.7in" y="1.65in" w="11.95in" h="4.6in" cols="3" rows="1" gap="0.25in">
              <stack col="1" fill="lt2" line="accent1" radius="6pt" pad="0.18in" gap="0.08in">
                <text size="14pt" bold="true" color="accent1">Our understanding</text>
                <text size="10pt" list="bullet"><p>Speed matters</p><p><run underline="true" link="https://example.com">Accessibility matters</run></p></text>
              </stack>
              <stack col="2" fill="lt2" line="accent2" radius="6pt" pad="0.18in">
                <text size="14pt" bold="true" color="accent2">Risks</text>
                <text size="10pt" list="bullet"><p>Security dependency</p><p>Data readiness</p></text>
              </stack>
              <stack col="3" fill="lt2" line="accent3" radius="6pt" pad="0.18in">
                <text size="14pt" bold="true" color="accent3">Plan</text>
                <text size="10pt" list="number"><p>Discover</p><p>Build</p><p>Launch</p></text>
              </stack>
            </grid>
          </slide>
          <slide layout="titleBody" flow="free">
            <text placeholder="title">Fee table</text>
            <table x="0.85in" y="1.65in" w="11.65in" h="3.2in" cols="4" header="true" colWidths="40%,20%,20%,20%" line="dk2">
              <row><cell>Position</cell><cell align="r">Daily rate</cell><cell align="r">Days</cell><cell align="r">Total</cell></row>
              <row><cell>Partner</cell><cell align="r">1000</cell><cell align="r">10</cell><cell align="r">10000</cell></row>
              <row><cell>Senior Consultant</cell><cell align="r">2000</cell><cell align="r">20</cell><cell align="r">40000</cell></row>
            </table>
          </slide>
          <slide layout="blank" flow="free">
            <text x="0.7in" y="0.65in" w="8in" h="0.6in" size="24pt" bold="true">Edge-anchored classification bar</text>
            <shape x="0.7in" y="1.4in" w="2.6in" h="0.45in" fill="none" line="accent1" text="Outline only"/>
            <stack dir="h" anchor="bottom" align="stretch" gap="0" h="0.42in">
              <shape fill="#E53935" text="HIGH RISK CONFIDENTIAL"/>
              <shape fill="#FB8C00" text="CONFIDENTIAL"/>
              <shape fill="#43A047" text="PUBLIC"/>
            </stack>
          </slide>
          <slide layout="blank" flow="free">
            <timeline x="0.7in" y="0.8in" w="11.9in" h="5.2in" periods="5" labels="W1,W2,W3,W4,W5">
              <task label="Discover" start="1" span="2" tone="accent2"/>
              <task label="Build" start="2" span="3" tone="accent1"/>
              <task label="Launch" start="5" span="1" tone="accent3"/>
              <milestone label="Kickoff" at="1"/>
              <milestone label="Go-live" at="5"/>
            </timeline>
          </slide>
        </deck>
        """

        output_path = Path("tests/fixtures/transpiled_deck.pptx")
        with tempfile.NamedTemporaryFile("w", suffix=".xml", delete=False) as handle:
            handle.write(dsl)
            source_path = Path(handle.name)
        try:
            result = transpile_deck(source_path, output_path)
        finally:
            source_path.unlink(missing_ok=True)

        presentation = PptxPresentation(str(result.pptx_path))
        self.assertEqual(len(presentation.slides), 4)
        self.assertTrue(all(len(slide.shapes) > 0 for slide in presentation.slides))
        self.assertTrue(output_path.exists())
        self.assertFalse([issue for issue in result.validation_issues if issue.severity == "error"])

        with ZipFile(output_path) as pptx:
            slide_xml = "\n".join(
                pptx.read(name).decode("utf-8")
                for name in pptx.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )
        self.assertIn("<a:buChar", slide_xml)
        self.assertIn("<a:buAutoNum", slide_xml)
        self.assertIn("<a:tbl>", slide_xml)
        self.assertIn("<a:noFill", slide_xml)
        self.assertIn('u="sng"', slide_xml)
        self.assertIn("<p:cxnSp>", slide_xml)
    
    def test_charts_demo_transpiles_opens_and_contains_chart_parts(self):
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            self.fail("python-pptx is required for transpiler integration tests")

        source_path = Path("tests/fixtures/golden/charts_demo.xml")
        output_path = Path("tests/fixtures/charts_demo.pptx")
        result = transpile_deck(source_path, output_path)

        presentation = PptxPresentation(str(result.pptx_path))
        self.assertEqual(len(presentation.slides), 6)
        self.assertFalse([issue for issue in result.validation_issues if issue.severity == "error"])

        CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        ns = dict(NSMAP, c=CHART_NS)

        with ZipFile(output_path) as pptx:
            names = pptx.namelist()

            # 1. chart parts + embedded workbooks exist, one per chart
            chart_parts = sorted(n for n in names if n.startswith("ppt/charts/chart") and n.endswith(".xml"))
            workbooks = sorted(n for n in names if n.startswith("xl/embeddings/") and n.endswith(".xlsx"))
            self.assertEqual(len(chart_parts), 6)
            self.assertEqual(len(workbooks), 6)

            # 2. correct plot element per chart type
            plot_elements = []
            for part in chart_parts:
                root = parse_xml(pptx.read(part).decode("utf-8"))
                plot_area = root.find("c:chart/c:plotArea", ns)
                kinds = [parse_xml_localname(child) for child in plot_area]
                plot_elements.append([k for k in kinds if k.endswith("Chart")])

            flat_plots = [p for sub in plot_elements for p in sub]
            self.assertIn("barChart", flat_plots)
            self.assertIn("lineChart", flat_plots)
            self.assertIn("pieChart", flat_plots)
            self.assertIn("areaChart", flat_plots)
            self.assertIn("scatterChart", flat_plots)

            # 3. cache == workbook for the first chart (the invariant)
            self._assert_cache_matches_workbook(pptx, chart_parts[0], workbooks, ns)

            # 4. finishes emit their distinct spPr signatures somewhere
            all_chart_xml = "\n".join(pptx.read(p).decode("utf-8") for p in chart_parts)
            self.assertIn("<a:gradFill", all_chart_xml)   # gradient-subtle slide

            # 5. no out-of-range luminance anywhere (repair-bug guard)
            luminance = []
            for name in names:
                if name.endswith(".xml"):
                    root = parse_xml(pptx.read(name).decode("utf-8"))
                    for transform in ("lumMod", "lumOff"):
                        luminance.extend(int(node.get("val")) for node in root.findall(f".//a:{transform}", ns))
            self.assertTrue(luminance)
            self.assertLessEqual(max(luminance), 100000)

        self.assertTrue(output_path.exists())

    def test_combo_chart_groups_axes_and_cache(self):
        from zipfile import ZipFile
        import io
        from openpyxl import load_workbook

        source_path = Path("tests/fixtures/golden/charts_demo.xml")
        output_path = Path("tests/fixtures/charts_demo.pptx")
        transpile_deck(source_path, output_path)

        CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        ns = dict(NSMAP, c=CHART_NS)

        with ZipFile(output_path) as pptx:
            names = pptx.namelist()
            chart_parts = sorted(n for n in names if n.startswith("ppt/charts/chart") and n.endswith(".xml"))

            # locate the combo part: the one plotArea with >1 chart-group element
            combo_part = None
            for part in chart_parts:
                root = parse_xml(pptx.read(part).decode("utf-8"))
                plot_area = root.find("c:chart/c:plotArea", ns)
                groups = [c for c in plot_area if parse_xml_localname(c).endswith("Chart")]
                if len(groups) > 1:
                    combo_part = part
                    combo_root = root
                    combo_plot_area = plot_area
                    combo_groups = groups
                    break
            self.assertIsNotNone(combo_part, "no combo chart part found")

            # 1. exactly one barChart (2 ser) + one lineChart (1 ser); total ser == input count
            by_kind = {parse_xml_localname(g): g for g in combo_groups}
            self.assertEqual(sorted(by_kind), ["barChart", "lineChart"])
            self.assertEqual(len(by_kind["barChart"].findall("c:ser", ns)), 2)
            self.assertEqual(len(by_kind["lineChart"].findall("c:ser", ns)), 1)
            total_ser = sum(len(g.findall("c:ser", ns)) for g in combo_groups)
            self.assertEqual(total_ser, 3)   # <-- the assertion that catches the doubling bug

            # 2. four distinct axis ids; every plot-group axId and every crossAx resolves
            axes = [a for a in combo_plot_area if parse_xml_localname(a).endswith("Ax")]
            axis_ids = [a.find("c:axId", ns).get("val") for a in axes]
            self.assertEqual(len(axis_ids), 4)
            self.assertEqual(len(set(axis_ids)), 4)   # distinct
            axis_id_set = set(axis_ids)
            for g in combo_groups:
                for ax in g.findall("c:axId", ns):
                    self.assertIn(ax.get("val"), axis_id_set)
            for a in axes:
                self.assertIn(a.find("c:crossAx", ns).get("val"), axis_id_set)

            # 3. secondary valAx: axPos=r + crosses=max ; secondary catAx: delete=1
            val_axes = [a for a in axes if parse_xml_localname(a) == "valAx"]
            cat_axes = [a for a in axes if parse_xml_localname(a) == "catAx"]
            self.assertTrue(any(
                a.find("c:axPos", ns).get("val") == "r"
                and a.find("c:crosses", ns) is not None
                and a.find("c:crosses", ns).get("val") == "max"
                for a in val_axes
            ), "no secondary value axis on the right crossing at max")
            self.assertTrue(any(
                a.find("c:delete", ns).get("val") == "1" for a in cat_axes
            ), "no hidden secondary category axis")

            # 4. cache == workbook for the combo chart (paired by the chart's own r:id)
            rels = parse_xml(pptx.read(f"ppt/charts/_rels/{Path(combo_part).name}.rels").decode("utf-8"))
            rel_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
            ext_rid = combo_root.find("c:externalData", ns).get(qn("r", "id"))
            target = None
            for rel in rels.findall("r:Relationship", rel_ns):
                if rel.get("Id") == ext_rid:
                    target = rel.get("Target")
                    break
            self.assertIsNotNone(target, "combo chart has no externalData relationship")
            wb_name = "xl/embeddings/" + Path(target).name
            cache_values = sorted(
                float(v.text) for v in combo_root.findall(".//c:val/c:numRef/c:numCache/c:pt/c:v", ns)
            )
            ws = load_workbook(io.BytesIO(pptx.read(wb_name))).active
            wb_values = []
            col = 2
            while ws.cell(row=2, column=col).value is not None:
                r = 2
                while ws.cell(row=r, column=col).value is not None:
                    wb_values.append(float(ws.cell(row=r, column=col).value))
                    r += 1
                col += 1
            self.assertEqual(cache_values, sorted(wb_values))

    def _assert_cache_matches_workbook(self, pptx, chart_part, workbooks, ns):
        """The chart's numCache values must equal the embedded workbook's cell values."""
        from openpyxl import load_workbook
        import io

        root = parse_xml(pptx.read(chart_part).decode("utf-8"))
        cache_values = [
            float(v.text)
            for v in root.findall(".//c:val/c:numRef/c:numCache/c:pt/c:v", ns)
        ]
        # load the first workbook and read its numeric cells (column B onward, rows 2+)
        wb = load_workbook(io.BytesIO(pptx.read(workbooks[0])))
        ws = wb.active
        wb_values = []
        col = 2  # B
        while ws.cell(row=2, column=col).value is not None:
            r = 2
            while ws.cell(row=r, column=col).value is not None:
                wb_values.append(float(ws.cell(row=r, column=col).value))
                r += 1
            col += 1
        self.assertEqual(sorted(cache_values), sorted(wb_values))


class EngineHardeningTests(unittest.TestCase):
    """Regression guards for fixes surfaced by the agent-DSL probe."""

    def test_xml_comments_do_not_crash_parser(self):  # #6
        parse_resolve(
            '<deck><!-- a note --><slide layout="blank" flow="stack">'
            '<text role="body">x</text></slide></deck>'
        )

    def test_chart_category_value_mismatch_rejected(self):  # #8
        with self.assertRaises(ValueError):
            parse_resolve(
                '<deck><slide layout="blank" flow="free">'
                '<chart type="column" x="1in" y="1in" w="6in" h="4in">'
                '<categories>A,B,C</categories>'
                '<series name="s"><point cat="A" value="1"/></series>'
                '</chart></slide></deck>'
            )

    def test_chart_matched_category_value_accepted(self):  # #8 must-not-false-positive
        parse_resolve(
            '<deck><slide layout="blank" flow="free">'
            '<chart type="column" x="1in" y="1in" w="6in" h="4in">'
            '<categories>A,B</categories>'
            '<series name="s"><point cat="A" value="1"/><point cat="B" value="2"/></series>'
            '</chart></slide></deck>'
        )

    def test_bad_theme_hex_rejected(self):  # #9
        ast = DSLParser().parse(
            '<deck><theme name="t" accent1="notacolor"/>'
            '<slide layout="blank" flow="stack"><text role="body">x</text></slide></deck>'
        )
        with self.assertRaises(Exception):  # DSLParseError
            _inline_theme_map(ast)

    def test_valid_theme_hex_accepted(self):  # #9 must-not-false-positive
        ast = DSLParser().parse(
            '<deck><theme name="t" accent1="2563EB"/>'
            '<slide layout="blank" flow="stack"><text role="body">x</text></slide></deck>'
        )
        self.assertEqual(_inline_theme_map(ast)['t']['colors']['accent1'], '2563EB')

    

if __name__ == "__main__":
    unittest.main()
